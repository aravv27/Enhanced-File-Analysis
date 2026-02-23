"""
Text extraction layer for AutoSorter.

Provides unified extract_text(filepath) dispatcher that routes to the
appropriate extractor based on file extension. Each extractor is wrapped
in try/except to gracefully handle corrupt or unreadable files.
"""

import json
import os

from src.config import get_file_type, SUPPORTED_EXTENSIONS, load_config
from src.logger import get_logger


def extract_text(filepath):
    """Extract text content from a file.
    
    Routes to the appropriate extractor based on file extension.
    Returns empty string on any extraction failure.
    
    Args:
        filepath: Absolute path to the file.
        
    Returns:
        str: Extracted text content, or empty string on failure.
    """
    logger = get_logger()
    ext = os.path.splitext(filepath)[1].lower()
    
    try:
        if ext == '.pdf':
            return _extract_pdf(filepath)
        elif ext == '.docx':
            return _extract_docx(filepath)
        elif ext == '.pptx':
            return _extract_pptx(filepath)
        elif ext in {'.jpg', '.jpeg', '.png'}:
            return _extract_image(filepath)
        elif ext == '.ipynb':
            return _extract_notebook(filepath)
        elif ext in {'.py', '.c', '.lex'}:
            return _extract_code(filepath)
        else:
            logger.warning(f"No extractor for extension: {ext}")
            return ""
    except Exception as e:
        logger.error(f"Extraction failed for {os.path.basename(filepath)}: {e}")
        return ""


def _extract_pdf(filepath):
    """Extract text from a PDF file.
    
    Attempts direct text extraction first. If result is empty or very short,
    falls back to OCR on page images. Limited to first N pages (from config).
    """
    import fitz  # PyMuPDF
    
    logger = get_logger()
    config = load_config()
    max_pages = config.get('ocr_max_pages', 5)
    
    doc = fitz.open(filepath)
    num_pages = min(len(doc), max_pages)
    
    # Attempt direct text extraction
    text_parts = []
    for page_num in range(num_pages):
        page = doc[page_num]
        text_parts.append(page.get_text())
    
    text = "\n".join(text_parts).strip()
    
    # If direct extraction yielded very little text, try OCR
    if len(text) < 50:
        logger.info(f"PDF direct extraction yielded little text, attempting OCR: {os.path.basename(filepath)}")
        text = _ocr_pdf_pages(doc, num_pages)
    
    doc.close()
    return text


def _ocr_pdf_pages(doc, num_pages):
    """OCR fallback for PDF pages that lack embedded text."""
    import pytesseract
    from PIL import Image
    import io
    
    text_parts = []
    for page_num in range(num_pages):
        page = doc[page_num]
        # Render page to image at 200 DPI
        pix = page.get_pixmap(dpi=200)
        img_data = pix.tobytes("png")
        img = Image.open(io.BytesIO(img_data))
        page_text = pytesseract.image_to_string(img)
        text_parts.append(page_text)
    
    return "\n".join(text_parts).strip()


def _extract_docx(filepath):
    """Extract full text from a DOCX file."""
    from docx import Document
    
    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _extract_pptx(filepath):
    """Extract text from all slides in a PPTX file."""
    from pptx import Presentation
    
    prs = Presentation(filepath)
    text_parts = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
    
    return "\n".join(text_parts)


def _extract_image(filepath):
    """Extract text from an image using OCR."""
    import pytesseract
    from PIL import Image
    
    img = Image.open(filepath)
    text = pytesseract.image_to_string(img)
    return text.strip()


def _extract_notebook(filepath):
    """Extract text from a Jupyter notebook (.ipynb).
    
    Reads source content from both code and markdown cells.
    """
    config = load_config()
    max_lines = config.get('code_max_lines', 500)
    
    with open(filepath, 'r', encoding='utf-8') as f:
        notebook = json.load(f)
    
    text_parts = []
    total_lines = 0
    
    cells = notebook.get('cells', [])
    for cell in cells:
        source = cell.get('source', [])
        if isinstance(source, list):
            lines = source
        else:
            lines = source.split('\n')
        
        for line in lines:
            if total_lines >= max_lines:
                break
            text_parts.append(line.rstrip('\n'))
            total_lines += 1
        
        if total_lines >= max_lines:
            break
    
    return "\n".join(text_parts)


def _extract_code(filepath):
    """Extract raw text from a code file (.py, .c, .lex).
    
    Reads up to the configured maximum number of lines.
    """
    config = load_config()
    max_lines = config.get('code_max_lines', 500)
    
    lines = []
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        for i, line in enumerate(f):
            if i >= max_lines:
                break
            lines.append(line.rstrip('\n'))
    
    return "\n".join(lines)
