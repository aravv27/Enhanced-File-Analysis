"""
AutoSorter Stress Test Suite.

Tests the full system end-to-end by launching the real app as a subprocess,
dropping N test files into a temp "Downloads" folder, and measuring how
fast the system detects, processes, and moves them.

Modes:
    --mock : Skips classification (instant), tests pure system throughput.
    --real : Uses actual embedding model, tests complete pipeline speed.

Usage:
    python -m tests.stress_test --mock          # System throughput only
    python -m tests.stress_test --real           # Full pipeline
    python -m tests.stress_test --mock --type py # Single file type
"""

import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
import time
from datetime import datetime

# ---------------------------------------------------------------------------
# File Generators — create minimal valid files of each type
# ---------------------------------------------------------------------------

def _gen_py(path):
    """Generate a minimal Python file."""
    with open(path, 'w', encoding='utf-8') as f:
        f.write(
            '# Auto-generated test file\n'
            'def hello():\n'
            '    print("Hello from AutoSorter stress test")\n'
            '\n'
            'class Node:\n'
            '    def __init__(self, data):\n'
            '        self.data = data\n'
            '        self.next = None\n'
            '\n'
            'if __name__ == "__main__":\n'
            '    hello()\n'
        )


def _gen_c(path):
    """Generate a minimal C file."""
    with open(path, 'w', encoding='utf-8') as f:
        f.write(
            '#include <stdio.h>\n'
            'int main() {\n'
            '    printf("Hello from AutoSorter stress test\\n");\n'
            '    return 0;\n'
            '}\n'
        )


def _gen_lex(path):
    """Generate a minimal Lex file."""
    with open(path, 'w', encoding='utf-8') as f:
        f.write(
            '%{\n'
            '#include <stdio.h>\n'
            '%}\n'
            '%%\n'
            '[a-zA-Z]+  printf("WORD ");\n'
            '[0-9]+     printf("NUM ");\n'
            '.          ;\n'
            '%%\n'
            'int main() { yylex(); return 0; }\n'
        )


def _gen_ipynb(path):
    """Generate a minimal Jupyter notebook."""
    nb = {
        "cells": [{
            "cell_type": "code",
            "source": ["print('Stress test notebook')\n", "x = [1, 2, 3]\n"],
            "metadata": {},
            "outputs": [],
            "execution_count": None,
        }],
        "metadata": {
            "kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
            "language_info": {"name": "python", "version": "3.10.0"},
        },
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(nb, f)


def _gen_docx(path):
    """Generate a minimal .docx file."""
    from docx import Document
    doc = Document()
    doc.add_paragraph(
        "Reinforcement learning and Q-learning are fundamental topics in machine learning. "
        "This document discusses Markov decision processes and policy optimization."
    )
    doc.save(path)


def _gen_pptx(path):
    """Generate a minimal .pptx file."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = "Stress Test Slide"
    slide.placeholders[1].text = (
        "public static void main(String[] args) {\n"
        "    System.out.println(\"Hello World\");\n"
        "}\n"
    )
    prs.save(path)


def _gen_pdf(path):
    """Generate a minimal 1-page PDF."""
    import fitz  # PyMuPDF
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text(
        (72, 72),
        "Neural networks and deep learning form the backbone of modern AI.\n"
        "Convolutional layers, backpropagation, and gradient descent are key concepts.",
        fontsize=12,
    )
    doc.save(path)
    doc.close()


def _gen_png(path):
    """Generate a minimal PNG with text."""
    from PIL import Image, ImageDraw
    img = Image.new('RGB', (200, 100), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.text((10, 30), "Stress test image", fill=(0, 0, 0))
    img.save(path)


def _gen_jpg(path):
    """Generate a minimal JPEG with text."""
    from PIL import Image, ImageDraw
    img = Image.new('RGB', (200, 100), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    draw.text((10, 30), "Stress test JPEG", fill=(0, 0, 0))
    img.save(path, format='JPEG')


# Map of file type -> generator function
GENERATORS = {
    '.py':    _gen_py,
    '.c':     _gen_c,
    '.lex':   _gen_lex,
    '.ipynb': _gen_ipynb,
    '.docx':  _gen_docx,
    '.pptx':  _gen_pptx,
    '.pdf':   _gen_pdf,
    '.png':   _gen_png,
    '.jpg':   _gen_jpg,
}

# N values to test
N_VALUES = [1, 5, 10, 50, 100, 250, 500, 1000, 2000, 2500, 5000, 10000]


# ---------------------------------------------------------------------------
# Test Config Generator
# ---------------------------------------------------------------------------

def create_test_config(source_dir, dest_dir, mock_mode, project_root):
    """Create a temporary config.json for the stress test.

    Args:
        source_dir: Temp directory to watch (fake Downloads).
        dest_dir: Temp directory for output (fake Subjects).
        mock_mode: If True, sets a flag the app can use to skip classification.
        project_root: Path to the project root.

    Returns:
        Path to the generated config file.
    """
    config = {
        "source_dir": source_dir,
        "destination_dir": dest_dir,
        "scan_existing_on_startup": True,
        "confidence_threshold": 0.50,
        "max_file_size_mb": 100,
        "worker_threads": 4,  # Use 4 threads for stress testing
        "model_name": "all-MiniLM-L6-v2",
        "watch_delay_seconds": 0,  # No delay for speed
        "ocr_max_pages": 1,
        "code_max_lines": 50,
        "log_max_bytes": 52428800,  # 50MB log for stress test
        "log_backup_count": 1,
        "ignored_extensions": [".crdownload", ".tmp", ".part", ".partial"],
        "mock_classifier": mock_mode,  # Custom flag for mock mode
    }

    config_path = os.path.join(project_root, 'config', 'config_stress_test.json')
    with open(config_path, 'w', encoding='utf-8') as f:
        json.dump(config, f, indent=2)

    return config_path


# ---------------------------------------------------------------------------
# File Generation
# ---------------------------------------------------------------------------

def generate_files(staging_dir, ext, n):
    """Generate N files of the given extension in the staging directory.

    Args:
        staging_dir: Directory to create files in.
        ext: File extension (e.g., '.py').
        n: Number of files to create.
    """
    gen_func = GENERATORS[ext]
    for i in range(n):
        fname = f"stress_test_{i:05d}{ext}"
        fpath = os.path.join(staging_dir, fname)
        gen_func(fpath)


# ---------------------------------------------------------------------------
# Core Test Runner
# ---------------------------------------------------------------------------

def wait_for_processing(processed_path, expected_count, timeout=600):
    """Wait until processed_files.json shows the expected count.

    Args:
        processed_path: Path to the processed_files.json.
        expected_count: Number of files we expect to be processed.
        timeout: Max seconds to wait.

    Returns:
        Tuple of (actual_count, timed_out).
    """
    start = time.time()
    last_count = 0
    stall_start = None
    stall_timeout = 60  # If no progress for 60s, consider it done

    while time.time() - start < timeout:
        try:
            if os.path.exists(processed_path):
                with open(processed_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                count = len(data)
            else:
                count = 0
        except (json.JSONDecodeError, OSError):
            count = last_count

        if count >= expected_count:
            return count, False

        # Detect stall
        if count > last_count:
            last_count = count
            stall_start = time.time()
        elif stall_start and (time.time() - stall_start > stall_timeout):
            # No progress for stall_timeout seconds, assume done
            return count, True

        time.sleep(0.5)

    return last_count, True


def run_single_test(ext, n, mock_mode, project_root):
    """Run a single stress test: N files of type ext.

    Launches the app as a subprocess, drops files, waits for processing,
    and measures throughput. Subprocess output goes to a log file for debugging.

    Args:
        ext: File extension.
        n: Number of files.
        mock_mode: Whether to use the mock classifier.
        project_root: Path to the project root.

    Returns:
        Dict with test results.
    """
    # Create temp directories
    temp_base = tempfile.mkdtemp(prefix='autosorter_stress_')
    source_dir = os.path.join(temp_base, 'downloads')
    dest_dir = os.path.join(temp_base, 'subjects')
    staging_dir = os.path.join(temp_base, 'staging')
    subprocess_log = os.path.join(temp_base, 'subprocess.log')
    processed_path = os.path.join(
        os.environ.get('LOCALAPPDATA', os.path.join(os.path.expanduser('~'), 'AppData', 'Local')),
        'AutoSorter', 'processed_files.json'
    )

    os.makedirs(source_dir)
    os.makedirs(dest_dir)
    os.makedirs(staging_dir)

    result = {
        'ext': ext,
        'n': n,
        'mode': 'mock' if mock_mode else 'real',
        'gen_time': 0,
        'total_time': 0,
        'process_time': 0,
        'processed_count': 0,
        'files_per_sec': 0,
        'timed_out': False,
        'error': None,
    }

    proc = None
    log_file = None

    try:
        # Step 1: Generate files in staging
        print(f"  Generating {n} {ext} files...", end=' ', flush=True)
        gen_start = time.time()
        generate_files(staging_dir, ext, n)
        result['gen_time'] = time.time() - gen_start
        print(f"done ({result['gen_time']:.2f}s)")

        # Step 2: Create test config
        config_path = create_test_config(source_dir, dest_dir, mock_mode, project_root)

        # Step 3: Clear processed files registry
        if os.path.exists(processed_path):
            os.remove(processed_path)

        # Step 4: Start the app as subprocess — output to log file
        log_file = open(subprocess_log, 'w', encoding='utf-8')
        cmd = [sys.executable, '-m', 'src.main', '--config', config_path]
        
        print(f"  Starting AutoSorter (mode={'mock' if mock_mode else 'real'})...", end=' ', flush=True)
        proc = subprocess.Popen(
            cmd,
            cwd=project_root,
            stdout=log_file,
            stderr=subprocess.STDOUT,  # Merge stderr into log
            creationflags=subprocess.CREATE_NEW_PROCESS_GROUP if os.name == 'nt' else 0,
        )

        # Wait for startup and verify the process is alive
        startup_wait = 3 if mock_mode else 10
        time.sleep(startup_wait)

        if proc.poll() is not None:
            # Process has already exited — read the log for the error
            log_file.close()
            log_file = None
            with open(subprocess_log, 'r', encoding='utf-8') as f:
                error_output = f.read()
            result['error'] = f"Subprocess exited (code={proc.returncode})"
            print(f"FAILED (exit code {proc.returncode})")
            print(f"  Subprocess output:\n{error_output[:500]}")
            return result

        print("running")

        # Step 5: Copy all files from staging to source (simulates download burst)
        print(f"  Dropping {n} files into source dir...", end=' ', flush=True)
        drop_start = time.time()
        for fname in os.listdir(staging_dir):
            shutil.copy2(
                os.path.join(staging_dir, fname),
                os.path.join(source_dir, fname),
            )
        drop_time = time.time() - drop_start
        print(f"done ({drop_time:.2f}s)")

        # Step 6: Wait for processing to complete
        print(f"  Waiting for {n} files to be processed...", flush=True)
        process_start = time.time()
        processed_count, timed_out = wait_for_processing(processed_path, n, timeout=600)
        process_time = time.time() - process_start

        result['total_time'] = drop_time + process_time
        result['process_time'] = process_time
        result['processed_count'] = processed_count
        result['timed_out'] = timed_out
        result['files_per_sec'] = processed_count / process_time if process_time > 0 else 0

        status = "TIMEOUT" if timed_out else "OK"
        print(
            f"  Result: {processed_count}/{n} files | "
            f"{process_time:.2f}s | "
            f"{result['files_per_sec']:.1f} files/s | "
            f"{status}"
        )

        # If timed out, show subprocess log tail for debugging
        if timed_out:
            log_file.flush()
            try:
                with open(subprocess_log, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
                    tail = ''.join(lines[-10:])
                    print(f"  Subprocess log (last 10 lines):\n{tail}")
            except Exception:
                pass

    except Exception as e:
        result['error'] = str(e)
        print(f"  ERROR: {e}")

    finally:
        # Step 7: Close log file
        if log_file and not log_file.closed:
            log_file.close()

        # Step 8: Kill the app
        if proc is not None:
            try:
                proc.terminate()
                proc.wait(timeout=10)
            except Exception:
                try:
                    proc.kill()
                except Exception:
                    pass

        # Step 9: Clean up temp directories
        try:
            shutil.rmtree(temp_base, ignore_errors=True)
        except Exception:
            pass

        # Clean up test config
        try:
            cfg = os.path.join(project_root, 'config', 'config_stress_test.json')
            if os.path.exists(cfg):
                os.remove(cfg)
        except Exception:
            pass

    return result


# ---------------------------------------------------------------------------
# Main Runner
# ---------------------------------------------------------------------------

def print_results_table(results):
    """Print a formatted results table."""
    print("\n" + "=" * 90)
    print("STRESS TEST RESULTS")
    print("=" * 90)

    # Group by extension
    by_ext = {}
    for r in results:
        by_ext.setdefault(r['ext'], []).append(r)

    for ext, ext_results in by_ext.items():
        print(f"\n--- {ext} ({ext_results[0]['mode']} mode) ---")
        print(f"{'N':>7} | {'Processed':>10} | {'Time (s)':>10} | {'Files/s':>10} | {'Status':>8}")
        print("-" * 60)
        for r in sorted(ext_results, key=lambda x: x['n']):
            status = "TIMEOUT" if r['timed_out'] else ("ERROR" if r['error'] else "OK")
            print(
                f"{r['n']:>7} | "
                f"{r['processed_count']:>10} | "
                f"{r['process_time']:>10.2f} | "
                f"{r['files_per_sec']:>10.1f} | "
                f"{status:>8}"
            )

    print("\n" + "=" * 90)


def save_results(results, project_root):
    """Save results to a JSON file."""
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    mode = results[0]['mode'] if results else 'unknown'
    out_path = os.path.join(project_root, 'tests', f'stress_results_{mode}_{ts}.json')
    with open(out_path, 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=2)
    print(f"\nResults saved to: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Graph Plotting
# ---------------------------------------------------------------------------

# Color palette for file types
TYPE_COLORS = {
    '.py': '#3572A5',    # Python blue
    '.c': '#555555',     # C gray
    '.lex': '#8B4513',   # Lex brown
    '.ipynb': '#F37626', # Jupyter orange
    '.docx': '#2B579A',  # Word blue
    '.pptx': '#D24726',  # PowerPoint red
    '.pdf': '#E44D26',   # PDF red-orange
    '.png': '#4CAF50',   # Image green
    '.jpg': '#66BB6A',   # JPEG green
}


def plot_graphs(results, project_root):
    """Generate 4 graph images (2 subplots each) from stress test results.

    Saves to tests/graphs/ directory.

    Args:
        results: List of result dicts from test runs.
        project_root: Path to the project root.
    """
    try:
        import matplotlib
        matplotlib.use('Agg')  # Non-interactive backend
        import matplotlib.pyplot as plt
        import matplotlib.ticker as ticker
    except ImportError:
        print("\n⚠ matplotlib not installed — skipping graph generation.")
        print("  Install with: pip install matplotlib")
        return

    graphs_dir = os.path.join(project_root, 'tests', 'graphs')
    os.makedirs(graphs_dir, exist_ok=True)

    mode = results[0]['mode'] if results else 'unknown'
    mode_label = "Mock (System Only)" if mode == 'mock' else "Real (Full Pipeline)"

    # Group by extension
    by_ext = {}
    for r in results:
        if r['error'] is None:
            by_ext.setdefault(r['ext'], []).append(r)

    # Sort each group by N
    for ext in by_ext:
        by_ext[ext].sort(key=lambda x: x['n'])

    plt.rcParams.update({
        'figure.facecolor': '#1a1a2e',
        'axes.facecolor': '#16213e',
        'axes.edgecolor': '#e0e0e0',
        'axes.labelcolor': '#e0e0e0',
        'text.color': '#e0e0e0',
        'xtick.color': '#e0e0e0',
        'ytick.color': '#e0e0e0',
        'grid.color': '#2a2a4a',
        'grid.alpha': 0.5,
        'legend.facecolor': '#16213e',
        'legend.edgecolor': '#4a4a6a',
        'legend.labelcolor': '#e0e0e0',
        'font.size': 10,
        'axes.titlesize': 13,
        'axes.titleweight': 'bold',
    })

    # ── Graph 1: Throughput + Total Time ──────────────────────────────────
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 6))
    fig.suptitle(f'AutoSorter Performance — {mode_label}', fontsize=16, fontweight='bold', y=0.98)

    for ext, data in by_ext.items():
        ns = [d['n'] for d in data]
        color = TYPE_COLORS.get(ext, '#ffffff')

        # Left: Throughput (files/sec)
        fps = [d['files_per_sec'] for d in data]
        ax1.plot(ns, fps, 'o-', label=ext, color=color, markersize=5, linewidth=2)

        # Right: Total time
        times = [d['process_time'] for d in data]
        ax2.plot(ns, times, 's-', label=ext, color=color, markersize=5, linewidth=2)

    ax1.set_xlabel('Number of Files (N)')
    ax1.set_ylabel('Throughput (files/sec)')
    ax1.set_title('Throughput vs File Count')
    ax1.set_xscale('log')
    ax1.legend(fontsize=8, ncol=2)
    ax1.grid(True, alpha=0.3)

    ax2.set_xlabel('Number of Files (N)')
    ax2.set_ylabel('Processing Time (seconds)')
    ax2.set_title('Total Processing Time vs File Count')
    ax2.set_xscale('log')
    ax2.set_yscale('log')
    ax2.legend(fontsize=8, ncol=2)
    ax2.grid(True, alpha=0.3)

    plt.tight_layout(rect=[0, 0, 1, 0.94])
    path1 = os.path.join(graphs_dir, 'throughput_and_time.png')
    fig.savefig(path1, dpi=150, bbox_inches='tight')
    plt.close(fig)
    print(f"  Saved: {path1}")

    # ── Graph 2: Per-file Time + Scaling Efficiency ───────────────────────
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 6))
    fig.suptitle(f'Per-File Performance — {mode_label}', fontsize=16, fontweight='bold', y=0.98)

    for ext, data in by_ext.items():
        ns = [d['n'] for d in data]
        color = TYPE_COLORS.get(ext, '#ffffff')

        # Left: Time per file
        tpf = [d['process_time'] / d['processed_count'] if d['processed_count'] > 0 else 0 for d in data]
        ax1.plot(ns, tpf, 'o-', label=ext, color=color, markersize=5, linewidth=2)

        # Right: Scaling efficiency (files/sec relative to N=smallest)
        if data[0]['files_per_sec'] > 0:
            baseline = data[0]['files_per_sec']
            efficiency = [d['files_per_sec'] / baseline * 100 for d in data]
            ax2.plot(ns, efficiency, 's-', label=ext, color=color, markersize=5, linewidth=2)

    ax1.set_xlabel('Number of Files (N)')
    ax1.set_ylabel('Time per File (seconds)')
    ax1.set_title('Average Processing Time per File')
    ax1.set_xscale('log')
    ax1.legend(fontsize=8, ncol=2)
    ax1.grid(True, alpha=0.3)

    ax2.set_xlabel('Number of Files (N)')
    ax2.set_ylabel('Relative Throughput (%)')
    ax2.set_title('Scaling Efficiency (100% = N=1 baseline)')
    ax2.set_xscale('log')
    ax2.axhline(y=100, color='#ff6b6b', linestyle='--', alpha=0.5, label='Ideal')
    ax2.legend(fontsize=8, ncol=2)
    ax2.grid(True, alpha=0.3)

    plt.tight_layout(rect=[0, 0, 1, 0.94])
    path2 = os.path.join(graphs_dir, 'per_file_and_scaling.png')
    fig.savefig(path2, dpi=150, bbox_inches='tight')
    plt.close(fig)
    print(f"  Saved: {path2}")

    # ── Graph 3: Type Comparison Bars ─────────────────────────────────────
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 6))
    fig.suptitle(f'File Type Comparison — {mode_label}', fontsize=16, fontweight='bold', y=0.98)

    # Find the largest N that all types have in common
    all_ns = set.intersection(*[set(d['n'] for d in data) for data in by_ext.values()]) if by_ext else set()
    if all_ns:
        max_n = max(all_ns)
        mid_n = sorted(all_ns)[len(all_ns) // 2]  # Median N

        # Left: Throughput bar at max N
        exts_sorted = sorted(by_ext.keys())
        throughputs = []
        colors = []
        for ext in exts_sorted:
            match = [d for d in by_ext[ext] if d['n'] == max_n]
            throughputs.append(match[0]['files_per_sec'] if match else 0)
            colors.append(TYPE_COLORS.get(ext, '#ffffff'))

        bars1 = ax1.bar(exts_sorted, throughputs, color=colors, edgecolor='#e0e0e0', linewidth=0.5)
        ax1.set_xlabel('File Type')
        ax1.set_ylabel('Throughput (files/sec)')
        ax1.set_title(f'Throughput by Type (N={max_n})')
        ax1.grid(True, alpha=0.3, axis='y')
        # Add value labels on bars
        for bar, val in zip(bars1, throughputs):
            ax1.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5,
                     f'{val:.1f}', ha='center', va='bottom', fontsize=9, color='#e0e0e0')

        # Right: Processing time bar at max N
        proc_times = []
        for ext in exts_sorted:
            match = [d for d in by_ext[ext] if d['n'] == max_n]
            proc_times.append(match[0]['process_time'] if match else 0)

        bars2 = ax2.bar(exts_sorted, proc_times, color=colors, edgecolor='#e0e0e0', linewidth=0.5)
        ax2.set_xlabel('File Type')
        ax2.set_ylabel('Processing Time (seconds)')
        ax2.set_title(f'Total Time by Type (N={max_n})')
        ax2.grid(True, alpha=0.3, axis='y')
        for bar, val in zip(bars2, proc_times):
            ax2.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                     f'{val:.1f}s', ha='center', va='bottom', fontsize=9, color='#e0e0e0')
    else:
        ax1.text(0.5, 0.5, 'No common N values\nacross types', transform=ax1.transAxes,
                 ha='center', va='center', fontsize=14)
        ax2.text(0.5, 0.5, 'No common N values\nacross types', transform=ax2.transAxes,
                 ha='center', va='center', fontsize=14)

    plt.tight_layout(rect=[0, 0, 1, 0.94])
    path3 = os.path.join(graphs_dir, 'type_comparison.png')
    fig.savefig(path3, dpi=150, bbox_inches='tight')
    plt.close(fig)
    print(f"  Saved: {path3}")

    # ── Graph 4: Generation Time + Success Rate ──────────────────────────
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 6))
    fig.suptitle(f'System Health — {mode_label}', fontsize=16, fontweight='bold', y=0.98)

    for ext, data in by_ext.items():
        ns = [d['n'] for d in data]
        color = TYPE_COLORS.get(ext, '#ffffff')

        # Left: File generation time
        gen_times = [d['gen_time'] for d in data]
        ax1.plot(ns, gen_times, 'o-', label=ext, color=color, markersize=5, linewidth=2)

        # Right: Success rate
        success = [d['processed_count'] / d['n'] * 100 if d['n'] > 0 else 0 for d in data]
        ax2.plot(ns, success, 's-', label=ext, color=color, markersize=5, linewidth=2)

    ax1.set_xlabel('Number of Files (N)')
    ax1.set_ylabel('Generation Time (seconds)')
    ax1.set_title('Test File Generation Time')
    ax1.set_xscale('log')
    ax1.legend(fontsize=8, ncol=2)
    ax1.grid(True, alpha=0.3)

    ax2.set_xlabel('Number of Files (N)')
    ax2.set_ylabel('Success Rate (%)')
    ax2.set_title('File Processing Success Rate')
    ax2.set_xscale('log')
    ax2.set_ylim(0, 110)
    ax2.axhline(y=100, color='#4CAF50', linestyle='--', alpha=0.5, label='100%')
    ax2.legend(fontsize=8, ncol=2)
    ax2.grid(True, alpha=0.3)

    plt.tight_layout(rect=[0, 0, 1, 0.94])
    path4 = os.path.join(graphs_dir, 'generation_and_success.png')
    fig.savefig(path4, dpi=150, bbox_inches='tight')
    plt.close(fig)
    print(f"  Saved: {path4}")

    print(f"\n✅ All graphs saved to: {graphs_dir}")


def main():
    parser = argparse.ArgumentParser(description='AutoSorter Stress Test')
    mode_group = parser.add_mutually_exclusive_group(required=True)
    mode_group.add_argument('--mock', action='store_true', help='Skip classification (test system throughput)')
    mode_group.add_argument('--real', action='store_true', help='Use real classifier (test full pipeline)')
    parser.add_argument('--type', type=str, default=None,
                        help=f"Test single file type. Options: {', '.join(GENERATORS.keys())}")
    parser.add_argument('--n', type=int, nargs='+', default=None,
                        help=f"Custom N values. Default: {N_VALUES}")
    parser.add_argument('--workers', type=int, default=4,
                        help="Number of worker threads (default: 4)")
    args = parser.parse_args()

    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    mock_mode = args.mock

    # Determine which file types to test
    if args.type:
        ext = args.type if args.type.startswith('.') else f'.{args.type}'
        if ext not in GENERATORS:
            print(f"Error: unsupported type '{ext}'. Options: {list(GENERATORS.keys())}")
            sys.exit(1)
        file_types = [ext]
    else:
        file_types = list(GENERATORS.keys())

    # Determine N values
    n_values = args.n or N_VALUES

    # For real mode, cap at reasonable N values
    if args.real and args.n is None:
        n_values = [v for v in N_VALUES if v <= 100]
        print("Real mode: capping N at 100 (use --n to override)")

    mode_str = "MOCK (no classification)" if mock_mode else "REAL (full pipeline)"
    print(f"\n{'=' * 60}")
    print(f"AutoSorter Stress Test")
    print(f"Mode:       {mode_str}")
    print(f"File types: {file_types}")
    print(f"N values:   {n_values}")
    print(f"Workers:    {args.workers}")
    print(f"{'=' * 60}\n")

    all_results = []

    for ext in file_types:
        print(f"\n{'─' * 50}")
        print(f"Testing: {ext}")
        print(f"{'─' * 50}")

        for n in n_values:
            print(f"\n[{ext} x {n}]")
            result = run_single_test(ext, n, mock_mode, project_root)
            all_results.append(result)

    # Print and save results
    print_results_table(all_results)
    results_path = save_results(all_results, project_root)

    # Generate graphs
    print("\nGenerating graphs...")
    plot_graphs(all_results, project_root)


if __name__ == '__main__':
    main()

